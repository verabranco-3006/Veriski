"""Generate Ops Review deck — Method 2 + Method 3 + framing slides.

Layout inspired by Gemini reference:
- Title slide
- Method 2 — Strategic Pillars + Success Measures
- Method 3 — Strategic Pillars + Resilience Metrics
- Operational Readiness & Obstacles + Next Steps

Sources:
- M2 vision & metrics: ProcessEngineering_Internal/v2mom_definition/method_2_strategic_document.md
- M3 problem statement: Confluence PE-STRAT-05
- Layout & palette: Gemini reference (uploads/Change Management.png, Incident Response.png, Apps Script)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT_PATH = r"C:\Users\vbr\Desktop\Veriski\Veriski\meetings\prep\2026-05-14-ops-review-m2-m3-slides.pptx"

# Palette — exact hex from Gemini Apps Script
DARK_BLUE = RGBColor(0x1E, 0x3A, 0x8A)        # #1e3a8a — titles
PRIMARY_BLUE = RGBColor(0x25, 0x63, 0xEB)     # #2563eb — accents, fills, icons
GREY_TEXT = RGBColor(0x4B, 0x55, 0x63)        # #4b5563 — body text
SOFT_GREY = RGBColor(0x7B, 0x87, 0x99)
LIGHT_BG = RGBColor(0xF1, 0xF4, 0xF9)
LIGHTER_BG = RGBColor(0xF8, 0xFA, 0xFC)
TRACK_BG = RGBColor(0xE2, 0xE8, 0xF0)
SUCCESS_GREEN = RGBColor(0x2F, 0x8F, 0x6F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xDD, 0xE3, 0xEC)


def set_run(run, text, size=12, bold=False, color=GREY_TEXT, font="Calibri"):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_section_label(slide, x, y, w, text, color=PRIMARY_BLUE, size=15):
    box = slide.shapes.add_textbox(x, y, w, Inches(0.35))
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, text, size=size, bold=True, color=color)


def add_title_section(slide, prs, title, subtitle):
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), prs.slide_width - Inches(1.0), Inches(0.7)
    )
    tf = title_box.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, title, size=28, bold=True, color=DARK_BLUE)

    sub_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.05), prs.slide_width - Inches(1.0), Inches(0.4)
    )
    stf = sub_box.text_frame
    stf.margin_left = 0
    stf.margin_top = 0
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    set_run(sr, subtitle, size=13, color=GREY_TEXT)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(1.6),
        prs.slide_width - Inches(1.0),
        Emu(9525),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_pillar(slide, x, y, w, h, glyph, name, body):
    icon_size = Inches(0.55)
    icon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, icon_size, icon_size)
    icon.adjustments[0] = 0.25
    icon.fill.solid()
    icon.fill.fore_color.rgb = PRIMARY_BLUE
    icon.line.fill.background()
    itf = icon.text_frame
    itf.margin_left = 0
    itf.margin_right = 0
    itf.margin_top = 0
    itf.margin_bottom = 0
    itf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ip = itf.paragraphs[0]
    ip.alignment = PP_ALIGN.CENTER
    ir = ip.add_run()
    set_run(ir, glyph, size=22, bold=False, color=WHITE, font="Segoe UI Symbol")

    text_x = x + icon_size + Inches(0.2)
    text_w = w - icon_size - Inches(0.2)
    name_box = slide.shapes.add_textbox(text_x, y - Inches(0.02), text_w, Inches(0.42))
    ntf = name_box.text_frame
    ntf.margin_left = 0
    ntf.margin_top = 0
    np = ntf.paragraphs[0]
    nr = np.add_run()
    set_run(nr, name, size=15, bold=True, color=DARK_BLUE)

    body_box = slide.shapes.add_textbox(text_x, y + Inches(0.38), text_w, h - Inches(0.4))
    btf = body_box.text_frame
    btf.word_wrap = True
    btf.margin_left = 0
    btf.margin_top = 0
    bp = btf.paragraphs[0]
    br = bp.add_run()
    set_run(br, body, size=10.5, color=GREY_TEXT)


def add_stat_tile(slide, x, y, w, h, value, label):
    tile = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    tile.adjustments[0] = 0.06
    tile.fill.solid()
    tile.fill.fore_color.rgb = LIGHT_BG
    tile.line.fill.background()
    tile.text_frame.text = ""

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x + Inches(0.08),
        y + Inches(0.18),
        Inches(0.07),
        h - Inches(0.36),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_BLUE
    bar.line.fill.background()

    val_box = slide.shapes.add_textbox(
        x + Inches(0.3), y + Inches(0.15), w - Inches(0.35), h - Inches(0.55)
    )
    vtf = val_box.text_frame
    vtf.margin_left = 0
    vtf.margin_top = 0
    vp = vtf.paragraphs[0]
    vr = vp.add_run()
    set_run(vr, value, size=30, bold=True, color=DARK_BLUE)

    lab_box = slide.shapes.add_textbox(
        x + Inches(0.3), y + h - Inches(0.4), w - Inches(0.35), Inches(0.28)
    )
    ltf = lab_box.text_frame
    ltf.margin_left = 0
    ltf.margin_top = 0
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    set_run(lr, label, size=9.5, bold=True, color=SOFT_GREY)


def add_progress_bar(slide, x, y, w, label, value, fill_pct):
    lab_box = slide.shapes.add_textbox(x, y, w - Inches(0.7), Inches(0.25))
    ltf = lab_box.text_frame
    ltf.margin_left = 0
    ltf.margin_top = 0
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    set_run(lr, label, size=10.5, bold=True, color=DARK_BLUE)

    val_box = slide.shapes.add_textbox(x + w - Inches(0.8), y, Inches(0.8), Inches(0.25))
    vtf = val_box.text_frame
    vtf.margin_left = 0
    vtf.margin_top = 0
    vp = vtf.paragraphs[0]
    vp.alignment = PP_ALIGN.RIGHT
    vr = vp.add_run()
    set_run(vr, value, size=10.5, bold=True, color=DARK_BLUE)

    track_y = y + Inches(0.28)
    track = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, track_y, w, Inches(0.13)
    )
    track.adjustments[0] = 0.5
    track.fill.solid()
    track.fill.fore_color.rgb = TRACK_BG
    track.line.fill.background()
    track.text_frame.text = ""

    fill_w = Emu(int(w * fill_pct / 100))
    if fill_w > 0:
        fill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, track_y, fill_w, Inches(0.13)
        )
        fill.adjustments[0] = 0.5
        fill.fill.solid()
        fill.fill.fore_color.rgb = PRIMARY_BLUE
        fill.line.fill.background()
        fill.text_frame.text = ""


def add_footer(slide, prs, text):
    footer_y = prs.slide_height - Inches(0.28)
    box = slide.shapes.add_textbox(
        Inches(0.5), footer_y, prs.slide_width - Inches(1.0), Inches(0.2)
    )
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    set_run(r, text, size=8.5, color=SOFT_GREY)


# ----------------------------- Slide builders -----------------------------

def build_title_slide(prs, *, title, subtitle, footnote):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title (left-aligned, large)
    tbox = slide.shapes.add_textbox(
        Inches(0.7), Inches(2.4), prs.slide_width - Inches(1.4), Inches(2.0)
    )
    tf = tbox.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    for i, line in enumerate(title.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        set_run(r, line, size=44, bold=True, color=DARK_BLUE)

    # Accent rule under title
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7),
        Inches(4.5),
        Inches(1.5),
        Emu(38100),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = PRIMARY_BLUE
    rule.line.fill.background()

    # Subtitle
    sbox = slide.shapes.add_textbox(
        Inches(0.7), Inches(4.7), prs.slide_width - Inches(1.4), Inches(0.8)
    )
    stf = sbox.text_frame
    stf.margin_left = 0
    stf.margin_top = 0
    for i, line in enumerate(subtitle.split("\n")):
        p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
        r = p.add_run()
        set_run(r, line, size=18, color=GREY_TEXT)

    # Footnote (bottom left)
    fnote = slide.shapes.add_textbox(
        Inches(0.7), prs.slide_height - Inches(0.7), prs.slide_width - Inches(1.4), Inches(0.4)
    )
    p = fnote.text_frame.paragraphs[0]
    r = p.add_run()
    set_run(r, footnote, size=11, color=SOFT_GREY)
    return slide


def add_outcomes_panel(slide, x, y, w, h, outcomes):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    panel.adjustments[0] = 0.05
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHTER_BG
    panel.line.color.rgb = LINE
    panel.line.width = Pt(0.75)
    panel.text_frame.text = ""

    tb = slide.shapes.add_textbox(
        x + Inches(0.2), y + Inches(0.12), w - Inches(0.4), h - Inches(0.24)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    for i, bullet in enumerate(outcomes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        r1 = p.add_run()
        set_run(r1, "✓  ", size=11, bold=True, color=SUCCESS_GREEN)
        r2 = p.add_run()
        set_run(r2, bullet, size=10.5, color=GREY_TEXT)


def build_method_slide(prs, *, title, subtitle, pillars, success_heading, tiles, outcomes, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_section(slide, prs, title, subtitle)

    # Left: Strategic Pillars
    left_x = Inches(0.5)
    pillar_w = Inches(5.2)
    add_section_label(slide, left_x, Inches(1.85), pillar_w, "Strategic Pillars")
    pillar_top = Inches(2.35)
    pillar_h = Inches(1.45)
    pillar_gap = Inches(0.18)
    for i, (glyph, name, body) in enumerate(pillars):
        y = pillar_top + (pillar_h + pillar_gap) * i
        add_pillar(slide, left_x, y, pillar_w, pillar_h, glyph, name, body)

    # Right top: Success metrics — always 2x2 grid of 4 tiles
    right_x = Inches(6.4)
    right_w = Inches(6.4)
    add_section_label(slide, right_x, Inches(1.85), right_w, success_heading)
    tiles_top = Inches(2.35)
    tile_gap = Inches(0.18)
    tile_w = (right_w - tile_gap) / 2
    tile_h = Inches(1.15)
    for i, (val, label) in enumerate(tiles):
        row = i // 2
        col = i % 2
        tx = right_x + (tile_w + tile_gap) * col
        ty = tiles_top + (tile_h + tile_gap) * row
        add_stat_tile(slide, tx, ty, tile_w, tile_h, val, label)
    tiles_bottom = tiles_top + tile_h * 2 + tile_gap

    # Right bottom: Outcomes panel ("What success looks like")
    outcomes_label_y = tiles_bottom + Inches(0.25)
    add_section_label(
        slide, right_x, outcomes_label_y, right_w, "What success looks like", color=SUCCESS_GREEN
    )
    panel_y = outcomes_label_y + Inches(0.4)
    panel_h = prs.slide_height - panel_y - Inches(0.5)
    add_outcomes_panel(slide, right_x, panel_y, right_w, panel_h, outcomes)

    add_footer(slide, prs, "Process Engineering · V2MOM Ops Review · 2026-05-14")
    slide.notes_slide.notes_text_frame.text = notes
    return slide


def build_obstacles_slide(prs, *, title, subtitle, obstacles, next_steps, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_section(slide, prs, title, subtitle)

    # Left: Key Obstacles
    left_x = Inches(0.5)
    col_w = Inches(5.9)
    add_section_label(slide, left_x, Inches(1.85), col_w, "Key Obstacles")
    item_top = Inches(2.35)
    item_h = Inches(1.45)
    item_gap = Inches(0.18)
    for i, (glyph, name, body) in enumerate(obstacles):
        y = item_top + (item_h + item_gap) * i
        add_pillar(slide, left_x, y, col_w, item_h, glyph, name, body)

    # Right: Next Steps as a styled panel
    right_x = Inches(6.7)
    right_w = Inches(6.1)
    add_section_label(slide, right_x, Inches(1.85), right_w, "Next Steps")
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, right_x, Inches(2.35), right_w, Inches(4.5)
    )
    panel.adjustments[0] = 0.04
    panel.fill.solid()
    panel.fill.fore_color.rgb = LIGHT_BG
    panel.line.fill.background()
    panel.text_frame.text = ""

    tb = slide.shapes.add_textbox(
        right_x + Inches(0.25), Inches(2.5), right_w - Inches(0.5), Inches(4.2)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    for i, (name, body) in enumerate(next_steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8 if i > 0 else 0)
        r1 = p.add_run()
        set_run(r1, "▸  ", size=14, bold=True, color=PRIMARY_BLUE)
        r2 = p.add_run()
        set_run(r2, name, size=13, bold=True, color=DARK_BLUE)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        r3 = p2.add_run()
        set_run(r3, "      " + body, size=11, color=GREY_TEXT)

    add_footer(slide, prs, "Process Engineering · V2MOM Ops Review · 2026-05-14")
    slide.notes_slide.notes_text_frame.text = notes
    return slide


# ----------------------------- Main -----------------------------

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # =================== SLIDE 1: Title ===================
    build_title_slide(
        prs,
        title="STRATEGIC OPERATIONS\nEFFICIENCY & RESILIENCE",
        subtitle="Methods 2 & 3 — V2MOM Briefing\nProcess Engineering · Ops Review · 2026-05-14",
        footnote="Speaker: Vera Branco · Team Lead, Process Engineering",
    )

    # =================== SLIDE 2: Method 2 ===================
    m2_pillars = [
        (
            "▣",
            "Pipeline Integrity",
            "Unify the Change Catalog and lock accountability across the 4 CM roles. Eliminate shadow changes and close the Invisible Risk Surface. Covers M2.1 + M2.2.",
        ),
        (
            "▲",
            "Standard Promotion",
            "Discovery ritual plus Lifecycle Automation. Mature Normals graduate to Standards; routine changes flow through automation, not SRE toil. Covers M2.3 + M2.6.",
        ),
        (
            "●",
            "CFR Correlation",
            "RFC ↔ incident link at asset-version level. CFR becomes a number Engineering Leadership trusts and acts on. Covers M2.4 + M2.5.",
        ),
    ]
    m2_tiles = [
        ("30%", "LEAD TIME REDUCTION"),
        ("60%", "STANDARD CHANGE VOL."),
        ("90%", "CFR CORRELATION ACCURACY"),
        ("< 5%", "UNCATALOGUED OPERATIONS"),
    ]
    m2_outcomes = [
        "Every change traces to a defined catalog article",
        "Each of the 4 CM roles has sharp scope and a refusal mechanism",
        "Standards flow through automation; SRE shifts to reliability owner",
        "CFR is a trusted signal Engineering Leadership acts on",
    ]
    m2_notes = (
        "Method 2 — speaker framing:\n"
        "• Vision: catalog-driven, role-accountable, automated, measurable end-to-end.\n"
        "• 3 pillars cover all 6 submethods:\n"
        "   - Pipeline Integrity = M2.1 Catalog + M2.2 Role Accountability\n"
        "   - Standard Promotion = M2.3 Standards + M2.6 Lifecycle Automation\n"
        "   - CFR Correlation = M2.4 Correlation + M2.5 Pipeline\n"
        "• Numbers on the slide are forward-looking targets. CONFIRM WITH LAURA before publication.\n"
        "• Supporting baseline numbers if asked: 731 RFCs → 992 mutations (2025, 1.35x); "
        "45h review+approval lead time; 70.59h ghost time in 'Implementing'.\n"
        "• Status today: M2.1 Foundations in flight (1,800 → 26 → 61 operation types); M2.5 docs complete.\n"
        "• Source: v2mom_definition/method_2_strategic_document.md"
    )
    build_method_slide(
        prs,
        title="METHOD 2: TRANSFORM CHANGE MANAGEMENT",
        subtitle="Catalog-driven, role-accountable, automated, and measurable end-to-end.",
        pillars=m2_pillars,
        success_heading="Success Measures",
        tiles=m2_tiles,
        outcomes=m2_outcomes,
        notes=m2_notes,
    )

    # =================== SLIDE 3: Method 3 ===================
    m3_pillars = [
        (
            "≡",
            "Classify Intake",
            "Define operating models for Service Incidents vs Support Cases. Make the alert-to-incident path explicit so false negatives can be declared. Covers M3.3 + M3.6.",
        ),
        (
            "⚡",
            "Govern Response",
            "Failure Management Governance and Status Page Governance. Visibility on every system-wide incident, end-to-end. Covers M3.1 + M3.2.",
        ),
        (
            "↻",
            "Close the Loop",
            "Team-led retrospectives and Problem Management. Recurring failures convert into structural fixes. Covers M3.4 + M3.5.",
        ),
    ]
    m3_tiles = [
        ("30m", "PUBLIC MTTA"),
        ("> 80%", "TRIAGE ACCURACY"),
        ("14d", "AVG. RCA LEAD TIME"),
        ("> 90%", "PROBLEM COVERAGE"),
    ]
    m3_outcomes = [
        "SI and Support Cases each have a defined operating model and OLAs",
        "Alerts have a defined path to declaration — no silent false negatives",
        "Status Page reflects every system-wide incident lifecycle",
        "Retros team-led, closed within target lead time; Problem Mgmt closes the loop",
    ]
    m3_notes = (
        "Method 3 — speaker framing:\n"
        "• Vision: auditable triage, transparent response, systematic learning.\n"
        "• 3 pillars regroup 4 underlying problems across 6 submethods:\n"
        "   - Classify Intake = M3.3 Triage + M3.6 Event Management\n"
        "     (Two new framings: SI vs Support Case operating models; declare false negatives.)\n"
        "   - Govern Response = M3.1 Failure Mgmt Governance + M3.2 Status Page Governance\n"
        "   - Close the Loop = M3.4 Problem Management + M3.5 Retrospective Transformation\n"
        "• Numbers on the slide are forward-looking targets. CONFIRM WITH INÊS before publication:\n"
        "   - 30m Public MTTA; > 80% Triage Accuracy; > 90% Problem Coverage; 14d Avg. RCA Lead Time (vs 72d baseline)\n"
        "   - 85% Alert Validation (Event Mgmt Filter); > 95% Incident Response Maturity Assessment\n"
        "• Supporting baseline numbers if asked: 44% of system-wide incidents had zero comms; "
        "7% followed the full lifecycle; 72-day avg retro lead time; 63% of retros reviewed by 2 people.\n"
        "• Source: Confluence PE-STRAT-05"
    )
    build_method_slide(
        prs,
        title="METHOD 3: INCIDENT RESPONSE RESILIENCE",
        subtitle="Auditable triage, transparent response, systematic learning.",
        pillars=m3_pillars,
        success_heading="Resilience Metrics",
        tiles=m3_tiles,
        outcomes=m3_outcomes,
        notes=m3_notes,
    )

    # =================== SLIDE 4: Obstacles + Next Steps ===================
    obstacles = [
        (
            "▣",
            "AI Skill Gap",
            "Mastering Process-as-Code platforms across the team. Capability ramp needed to operate the new motion at scale.",
        ),
        (
            "⚡",
            "Data Reliability",
            "Trust in automated correlation logic. Data quality has to land before CFR can be acted on by Engineering Leadership.",
        ),
        (
            "↻",
            "Cross-Org Dependencies",
            "SRE, GS, and Engineering alignment on operating models. Decision rights need to be settled with sponsors early.",
        ),
    ]
    next_steps = [
        ("Confirm targets with method owners",
         "Laura validates M2 measures; Inês validates M3 measures before any external publication."),
        ("Sequence dependencies",
         "M2: Foundations first. M3: Classify Intake first (M3.3 + M3.6 unblocks the rest)."),
        ("Stakeholder alignment",
         "VS Leader pre-reads on operating models for SI vs Support Case and Event Mgmt declaration."),
        ("2027 Strategic Metrics Blueprint",
         "Finalize the measurement layer that ties M2 and M3 outcomes to leadership decision-making."),
    ]
    obstacles_notes = (
        "Obstacles & Next Steps — speaker framing:\n"
        "• Obstacles are pulled from Gemini's deck framing. CONFIRM whether they reflect your view "
        "before delivering — especially the 'AI Skill Gap' framing.\n"
        "• Next steps tie the deck to action: target confirmation, sequencing, stakeholder alignment, "
        "and the 2027 metrics blueprint.\n"
        "• Land: this is a strategy briefing, the next concrete step is target sign-off with Laura and Inês."
    )
    build_obstacles_slide(
        prs,
        title="OPERATIONAL READINESS & OBSTACLES",
        subtitle="What stands between the strategy and execution, and what we do next.",
        obstacles=obstacles,
        next_steps=next_steps,
        notes=obstacles_notes,
    )

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
